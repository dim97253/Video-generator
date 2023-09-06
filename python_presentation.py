from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import gpt_request as gpt
import image_generator as imagen
from PIL import Image

def set_slide_background(slide, image_path):
    left = top = 0
    pic = slide.shapes.add_picture(image_path, left, top, width=prs.slide_width, height=prs.slide_height)

    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

print("Enter topic of the presentation")
topic = input()

slides_data = gpt.request_presentation(topic, 5)

# Create a new presentation object
prs = Presentation()
# Define slide layout. 0 is the layout for a title slide
slide_layout = prs.slide_layouts[0] 
# Add a slide with the defined layout
slide = prs.slides.add_slide(slide_layout)
# Select title and subtitle placeholders 
title = slide.shapes.title
subtitle = slide.placeholders[1]  # Using index for placeholder
# Set text for title and subtitle
title.text = topic
subtitle.text = ""

for data in slides_data:
    slide_layout = prs.slide_layouts[3] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]       
    subtitle.text = data['text']
    text_frame = subtitle.text_frame
    text_frame.fit_text()
    title.text = data['title']
    
    img = imagen.get_image(data['image'])
    img_path = 'temp.png'
    with open(img_path, "wb") as f:
        f.write(img)

    left = slide.placeholders[2].left
    top = slide.placeholders[2].top
    width = slide.placeholders[2].width
    height = slide.placeholders[2].height

    slide.shapes.add_picture(img_path, left, top, width, height)

    # Use Pillow to process the image transparency
    img = Image.open(img_path).convert("RGBA")
    transparent_img = Image.new("RGBA", img.size, (255, 255, 255, 0))
    blended_img = Image.blend(transparent_img, img, alpha=0.5)
    blended_img_path = 'transparent_image.png'
    blended_img.save(blended_img_path, "PNG")

    set_slide_background(slide, blended_img_path)


# Save the presentation
prs.save('test_presentation.pptx')